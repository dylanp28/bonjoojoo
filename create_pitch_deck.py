#!/usr/bin/env python3
"""
bonjoojoo Investor Pitch Deck Generator
Creates a professional PowerPoint presentation for lab-grown diamond luxury brand
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_pitch_deck():
    # Create presentation
    prs = Presentation()
    
    # Define colors (elegant jewelry brand palette)
    primary_color = RGBColor(26, 26, 26)      # Rich black
    accent_color = RGBColor(184, 134, 11)     # Gold accent  
    secondary_color = RGBColor(64, 64, 64)    # Medium gray
    light_gray = RGBColor(248, 248, 248)      # Light background
    
    def add_title_slide(title, subtitle=""):
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = secondary_color
        
        return slide
    
    def add_content_slide(title, content_points):
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = primary_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        tf = content_shape.text_frame
        tf.clear()
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = secondary_color
            p.level = 0
            p.space_after = Pt(8)
        
        return slide

    def add_two_column_slide(title, left_content, right_content):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = primary_color
        title_frame.paragraphs[0].font.bold = True
        
        # Left column
        left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5.5))
        left_frame = left_shape.text_frame
        left_frame.clear()
        
        for i, point in enumerate(left_content):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(14)
            p.font.color.rgb = secondary_color
            p.space_after = Pt(6)
        
        # Right column
        right_shape = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.5), Inches(5.5))
        right_frame = right_shape.text_frame
        right_frame.clear()
        
        for i, point in enumerate(right_content):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(14)
            p.font.color.rgb = secondary_color
            p.space_after = Pt(6)
        
        return slide

    # Slide 1: Title Slide
    add_title_slide("bonjoojoo", "Conscious Luxury • Lab-Grown Diamonds • Next Generation Jewelry")
    
    # Slide 2: Company Overview
    add_content_slide("Company Overview", [
        "• Premium lab-grown diamond jewelry brand",
        "• Sustainable luxury positioned for millennials & Gen Z",
        "• Direct-to-consumer model with concierge-level service",
        "• Curated collections vs. overwhelming catalogs",
        "• Founded on principles of ethics, quality, and transparency",
        "",
        "💎 Conscious luxury. Infinite possibilities."
    ])
    
    # Slide 3: The Problem
    add_content_slide("The Problem: Traditional Diamond Industry", [
        "Environmental Impact:",
        "  • Diamond mining displaces 250 tons of earth per carat",
        "  • Carbon footprint 57,000x higher than lab-grown",
        "",
        "Ethical Concerns:",
        "  • Conflict diamonds still exist in supply chains",
        "  • Worker exploitation in mining operations",
        "",
        "Consumer Pain Points:",
        "  • Overpriced due to De Beers monopoly legacy",
        "  • Overwhelming choices without guidance",
        "  • Declining engagement ring budgets (-13% since 2021)"
    ])
    
    # Slide 4: Our Solution
    add_content_slide("Our Solution: Lab-Grown Excellence", [
        "Superior Value:",
        "  • 60-80% cost savings vs. natural diamonds",
        "  • Identical chemical, physical & optical properties",
        "  • IGI/GIA certified quality assurance",
        "",
        "Differentiated Experience:",
        "  • Curated collections with storytelling",
        "  • Personal jewelry consultations",
        "  • Educational content and transparency",
        "  • 30-day returns, lifetime warranty",
        "",
        "Technology Advantage: CVD process for perfect diamonds"
    ])
    
    # Slide 5: Market Opportunity
    add_content_slide("Explosive Market Growth", [
        "Global Lab-Grown Diamond Market:",
        "  • 2024: $29.7 billion (current)",
        "  • 2034: $91.8 billion (forecast)",
        "  • CAGR: 14.15% (2025-2034)",
        "",
        "US Market Specifics:",
        "  • Lab-grown = 47.4% of engagement ring market",
        "  • Average price: $5,200 (down from $6,000)",
        "  • Millennials/Gen Z: 78% prefer ethical alternatives",
        "",
        "Target Market Size: $850M serviceable, $2.1B total"
    ])
    
    # Slide 6: SWOT Analysis (Two Column)
    add_two_column_slide("SWOT Analysis", 
        [
            "STRENGTHS:",
            "• First-mover advantage in curated luxury",
            "• 65-70% gross margins",
            "• Founder jewelry industry expertise", 
            "• Sustainable positioning",
            "• Scalable D2C model",
            "",
            "OPPORTUNITIES:",
            "• Rapid market growth (14.15% CAGR)",
            "• Generational shift to ethical consumption",
            "• International expansion potential",
            "• Adjacent categories (watches, men's)",
            "• B2B partnerships"
        ],
        [
            "WEAKNESSES:",
            "• New brand requiring education",
            "• Limited brand recognition",
            "• Supplier relationship dependence",
            "• Higher initial CAC",
            "",
            "THREATS:",
            "• Established players entering space",
            "• Potential price commoditization", 
            "• Natural diamond marketing pushback",
            "• Economic downturn affecting luxury"
        ]
    )
    
    # Slide 7: Competitive Landscape
    add_content_slide("Competitive Landscape", [
        "Brilliant Earth (NASDAQ:BRLT):",
        "  • Market cap: ~$450M, Revenue: $400M+",
        "  • Weakness: Higher prices, less personal service",
        "",
        "Clean Origin:",
        "  • Revenue: ~$50M, Value-focused positioning",
        "  • Weakness: Limited brand story, basic designs",
        "",
        "Our Competitive Advantages:",
        "  • Curated luxury positioning",
        "  • Personal consultation model",
        "  • Story-driven brand experience",
        "  • Superior customer service"
    ])
    
    # Slide 8: Business Model
    add_content_slide("Business Model & Unit Economics", [
        "Revenue Streams:",
        "  • Engagement rings (65%) • Fine jewelry (25%) • Accessories (10%)",
        "",
        "Unit Economics:",
        "  • Average order value: $2,750",
        "  • Gross margin: 67%", 
        "  • Customer acquisition cost: $125",
        "  • Customer lifetime value: $1,850",
        "  • LTV/CAC ratio: 14.8x",
        "",
        "Pricing: Premium positioning 20-30% above value players,",
        "10-20% below luxury incumbents"
    ])
    
    # Slide 9: Target Market
    add_content_slide("Target Customer Personas", [
        "Sustainable Sarah (26-32) - 40% of market:",
        "  • Income: $75K-$120K",
        "  • Values: Ethics, sustainability, authenticity",
        "  • Spending: $1,800-$3,500 on engagement rings",
        "",
        "Luxury Lauren (30-40) - 35% of market:",
        "  • Income: $120K-$250K",
        "  • Values: Quality, exclusivity, craftsmanship",
        "  • Spending: $3,500-$8,000 on fine jewelry",
        "",
        "Modern Michael (28-35) - 25% of market:",
        "  • Income: $80K-$150K, Primary ring researcher",
        "  • Values: Innovation, smart purchasing"
    ])
    
    # Slide 10: Financial Projections
    add_content_slide("Financial Projections", [
        "Year 1 (2025): $1.2M revenue, 450 customers",
        "  • Gross margin: 65% • Operating margin: -15%",
        "",
        "Year 2 (2026): $3.8M revenue, 1,380 customers",
        "  • Gross margin: 67% • Operating margin: 5%",
        "",
        "Year 3 (2027): $9.2M revenue, 3,350 customers", 
        "  • Gross margin: 69% • Operating margin: 18%",
        "",
        "Year 5 (2029): $28.5M revenue, 10,350 customers",
        "  • CAGR: 125% (Years 1-3)",
        "  • International revenue: 25% by Year 5"
    ])
    
    # Slide 11: Technology Platform
    add_content_slide("Technology & Operations", [
        "E-commerce Platform:",
        "  • Next.js/React frontend for performance",
        "  • Shopify Plus backend for scalability",
        "  • Mobile-optimized (60% of traffic)",
        "",
        "Competitive Tech Advantages:",
        "  • Virtual try-on using AR technology",
        "  • Diamond visualization and comparison tools",
        "  • Real-time chat with jewelry experts",
        "",
        "Supply Chain:",
        "  • Direct relationships with lab-grown producers",
        "  • IGI-certified suppliers • Same-day shipping"
    ])
    
    # Slide 12: Growth Strategy
    add_content_slide("Growth Strategy", [
        "Phase 1: Market Establishment (Years 1-2)",
        "  • Build brand awareness in engagement rings",
        "  • Optimize unit economics and operations",
        "",
        "Phase 2: Category Expansion (Years 2-4)", 
        "  • Expand into fine jewelry categories",
        "  • Launch men's jewelry and watches",
        "",
        "Phase 3: Market Dominance (Years 4-5)",
        "  • International expansion (Canada, UK, Australia)",
        "  • Flagship showrooms in major markets",
        "",
        "Growth Levers: Product expansion, geographic growth,",
        "B2B partnerships, LTV optimization"
    ])
    
    # Slide 13: Investment Ask
    add_content_slide("Investment Ask: $2.5M Series A", [
        "Use of Funds:",
        "  • Inventory & Working Capital: 40% ($1.0M)",
        "  • Marketing & Customer Acquisition: 35% ($875K)", 
        "  • Technology Development: 15% ($375K)",
        "  • Team Expansion: 10% ($250K)",
        "",
        "Key Milestones:",
        "  • Achieve $3.8M revenue in Year 2",
        "  • Reach positive operating margin",
        "  • Launch 3 new product categories",
        "  • Expand to 2 international markets",
        "",
        "Return Potential: $85M-$140M valuation at exit"
    ])
    
    # Slide 14: Team
    add_content_slide("Team & Leadership", [
        "Founder & CEO:",
        "  • 10+ years family jewelry industry experience",
        "  • UC Berkeley Data Science & Economics",
        "  • Deep understanding of diamond supply chains",
        "",
        "Advisory Board:",
        "  • Industry veteran from Brilliant Earth/Blue Nile",
        "  • E-commerce scaling expert (former Warby Parker)",
        "  • Sustainable business consultant",
        "",
        "Key Hires with Funding:",
        "  • VP Marketing • Head of Operations",
        "  • Senior Designer • Customer Success Manager"
    ])
    
    # Slide 15: Thank You
    add_title_slide("Thank You", "Questions & Discussion\n\nContact: hello@bonjoojoo.com")
    
    # Save presentation
    output_path = os.path.expanduser("~/.openclaw/workspace/bonjoojoo/bonjoojoo-pitch-deck.pptx")
    prs.save(output_path)
    
    print(f"✅ Pitch deck created successfully!")
    print(f"📄 File saved: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    create_pitch_deck()